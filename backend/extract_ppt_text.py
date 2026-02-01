import sys
import json
import os

from pptx import Presentation
import pdfplumber
from docx import Document


def extract_pptx(ppt_path):
    prs = Presentation(ppt_path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)


def extract_pdf(pdf_path):
    text = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)

    return "\n".join(text)


def extract_docx(docx_path):
    doc = Document(docx_path)
    return "\n".join(p.text for p in doc.paragraphs)


def extract_txt(txt_path):
    with open(txt_path, "r", encoding="utf-8") as f:
        return f.read()


def extract_text(file_path):
    extension = os.path.splitext(file_path)[1].lower()

    if extension == ".pptx":
        return extract_pptx(file_path)
    elif extension == ".pdf":
        return extract_pdf(file_path)
    elif extension == ".docx":
        return extract_docx(file_path)
    elif extension == ".txt":
        return extract_txt(file_path)
    else:
        raise ValueError("Unsupported file format")


if __name__ == "__main__":
    file_path = sys.argv[1]

    try:
        extracted_text = extract_text(file_path)
        print(json.dumps({"text": extracted_text}))
    except Exception as e:
        print(json.dumps({"error": str(e)}))
